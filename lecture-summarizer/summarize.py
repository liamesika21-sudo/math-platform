#!/usr/bin/env python3
"""
Lecture Summarizer — Automated Pipeline
========================================
1. Reads PDF/PPTX files from ./generate/
2. Extracts text content
3. Sends to Claude API with the system prompt
4. Receives HTML summary
5. Converts HTML to PDF and saves to ./results/
6. Moves source files to ./done/
"""

import os
import sys
import shutil
import glob
import time
import json
import anthropic

# ── Config ───────────────────────────────────────────────────────────────────

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
GENERATE_DIR = os.path.join(BASE_DIR, "generate")
RESULTS_DIR = os.path.join(BASE_DIR, "results")
DONE_DIR = os.path.join(BASE_DIR, "done")
PROMPT_FILE = os.path.join(BASE_DIR, "SYSTEM_PROMPT.md")

# Create directories if they don't exist
for d in [GENERATE_DIR, RESULTS_DIR, DONE_DIR]:
    os.makedirs(d, exist_ok=True)


# ── Step 1: Extract text from PDF ────────────────────────────────────────────

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract all text from a PDF file."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        text_parts = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                text_parts.append(f"=== SLIDE {i + 1} ===\n{text.strip()}")
        doc.close()
        return "\n\n".join(text_parts)
    except ImportError:
        # Fallback to pdf-parse via subprocess
        try:
            import subprocess
            result = subprocess.run(
                ["python3", "-c", f"""
import fitz
doc = fitz.open("{pdf_path}")
for i, page in enumerate(doc):
    text = page.get_text()
    if text.strip():
        print(f"=== SLIDE {{i+1}} ===")
        print(text.strip())
        print()
doc.close()
"""],
                capture_output=True, text=True
            )
            return result.stdout
        except Exception:
            print(f"  ERROR: Could not extract text from {pdf_path}")
            print("  Install PyMuPDF: pip install PyMuPDF")
            return ""


def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract all text from a PPTX file."""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text_parts = []
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            slide_texts.append(txt)
            if slide_texts:
                text_parts.append(f"=== SLIDE {i + 1} ===\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except ImportError:
        print(f"  ERROR: Could not extract text from {pptx_path}")
        print("  Install python-pptx: pip install python-pptx")
        return ""


def extract_text(file_path: str) -> str:
    """Extract text from PDF or PPTX."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    else:
        print(f"  SKIP: Unsupported file type: {ext}")
        return ""


# ── Step 2: Parse filename for metadata ──────────────────────────────────────

def parse_filename(filename: str) -> dict:
    """
    Try to extract metadata from filename.
    Expected patterns:
      - "lecture1.pdf" / "lecture2.pdf"
      - "Recitation_2.pdf"
      - "הרצאה 1 - אחרי 2.pdf"
      - Custom: "DS-Lecture-3-Sorting.pdf"

    Returns dict with: file_type, number, title_hint
    """
    name = os.path.splitext(filename)[0].lower()

    file_type = "lecture"
    number = 1
    title_hint = ""

    if "recit" in name or "תרגול" in name or "tutorial" in name:
        file_type = "recitation"
    elif "הרצאה" in name or "lecture" in name:
        file_type = "lecture"

    # Try to find a number
    import re
    nums = re.findall(r'\d+', name)
    if nums:
        number = int(nums[0])

    return {
        "file_type": file_type,
        "number": number,
        "title_hint": title_hint,
    }


# ── Step 3: Call Claude API ──────────────────────────────────────────────────

def load_system_prompt() -> str:
    """Load the system prompt from SYSTEM_PROMPT.md."""
    with open(PROMPT_FILE, "r", encoding="utf-8") as f:
        return f.read()


def generate_html_summary(
    extracted_text: str,
    file_type: str,
    number: int,
    course_name: str = "מבני נתונים",
    university: str = "אוניברסיטת רייכמן",
) -> str:
    """Send extracted text to Claude and get HTML summary back."""

    system_prompt = load_system_prompt()

    file_type_heb = "הרצאה" if file_type == "lecture" else "תרגול"

    user_message = f"""Please create a comprehensive HTML summary for the following {file_type}.

**Course:** {course_name}
**University:** {university}
**Type:** {file_type_heb} ({file_type})
**Number:** {number}

**Important:**
- Analyze the content and determine the lecture title yourself from the content.
- Determine an appropriate English subtitle.
- Cover ALL slides — do not skip anything educational.
- Use ALL the box types appropriately (def-box, thm-box, proof-box, example-box, note-box, code-box, question-box, solution-box).
- Include all proofs step by step.
- Include all pseudocode in code-box with LTR direction.
- Include all mathematical formulas with KaTeX notation.

Here is the full extracted content from all slides:

---

{extracted_text}

---

Output ONLY the complete HTML file. No explanations, no markdown code blocks — just the raw HTML starting with <!DOCTYPE html>."""

    client = anthropic.Anthropic()  # Uses ANTHROPIC_API_KEY env var

    print("  Calling Claude API (this may take 1-2 minutes)...")

    message = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=16000,
        system=system_prompt,
        messages=[{"role": "user", "content": user_message}],
    )

    html_content = message.content[0].text

    # Clean up — remove markdown code fences if Claude wrapped it
    if html_content.startswith("```"):
        lines = html_content.split("\n")
        # Remove first line (```html) and last line (```)
        if lines[-1].strip() == "```":
            lines = lines[1:-1]
        elif lines[0].startswith("```"):
            lines = lines[1:]
        html_content = "\n".join(lines)

    return html_content


# ── Step 4: Convert HTML to PDF ──────────────────────────────────────────────

def html_to_pdf(html_path: str, pdf_path: str) -> bool:
    """Convert HTML file to PDF using available tool."""
    import subprocess

    # Try playwright (best quality)
    try:
        result = subprocess.run(
            ["python3", "-c", f"""
import asyncio
from playwright.async_api import async_playwright

async def main():
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page()
        await page.goto('file://{os.path.abspath(html_path)}')
        await page.wait_for_timeout(2000)  # Wait for KaTeX to render
        await page.pdf(path='{os.path.abspath(pdf_path)}', format='A4', print_background=True)
        await browser.close()

asyncio.run(main())
"""],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            return True
        print(f"  Playwright failed: {result.stderr[:200]}")
    except Exception as e:
        print(f"  Playwright not available: {e}")

    # Try wkhtmltopdf
    try:
        result = subprocess.run(
            ["wkhtmltopdf", "--enable-local-file-access", "--encoding", "UTF-8",
             html_path, pdf_path],
            capture_output=True, text=True, timeout=60
        )
        if result.returncode == 0:
            return True
        print(f"  wkhtmltopdf failed: {result.stderr[:200]}")
    except FileNotFoundError:
        print("  wkhtmltopdf not found")
    except Exception as e:
        print(f"  wkhtmltopdf error: {e}")

    # Try Chrome headless
    for chrome in [
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "google-chrome", "chromium-browser", "chromium"
    ]:
        try:
            result = subprocess.run(
                [chrome, "--headless", "--disable-gpu", "--no-sandbox",
                 f"--print-to-pdf={os.path.abspath(pdf_path)}",
                 f"file://{os.path.abspath(html_path)}"],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                return True
        except (FileNotFoundError, Exception):
            continue

    print("  WARNING: Could not convert to PDF. HTML file saved to results/ instead.")
    return False


# ── Main Pipeline ────────────────────────────────────────────────────────────

def process_file(file_path: str, course_name: str, university: str):
    """Process a single file through the full pipeline."""
    filename = os.path.basename(file_path)
    print(f"\n{'='*60}")
    print(f"Processing: {filename}")
    print(f"{'='*60}")

    # 1. Extract text
    print("  [1/4] Extracting text...")
    text = extract_text(file_path)
    if not text:
        print(f"  FAILED: No text extracted from {filename}")
        return False

    slide_count = text.count("=== SLIDE")
    print(f"  Extracted {slide_count} slides")

    # 2. Parse metadata
    meta = parse_filename(filename)
    print(f"  Type: {meta['file_type']}, Number: {meta['number']}")

    # 3. Generate HTML via Claude
    print("  [2/4] Generating HTML summary via Claude API...")
    html = generate_html_summary(
        extracted_text=text,
        file_type=meta["file_type"],
        number=meta["number"],
        course_name=course_name,
        university=university,
    )

    if not html or "<!DOCTYPE" not in html.lower():
        print("  FAILED: Claude did not return valid HTML")
        return False

    # 4. Save HTML
    base_name = os.path.splitext(filename)[0]
    html_filename = f"{base_name}-summary.html"
    html_path = os.path.join(RESULTS_DIR, html_filename)

    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html)
    print(f"  [3/4] HTML saved: results/{html_filename}")

    # 5. Convert to PDF
    print("  [4/4] Converting to PDF...")
    pdf_filename = f"{base_name}-summary.pdf"
    pdf_path = os.path.join(RESULTS_DIR, pdf_filename)

    if html_to_pdf(html_path, pdf_path):
        print(f"  PDF saved: results/{pdf_filename}")
    else:
        # Keep HTML as fallback
        print(f"  HTML kept as fallback: results/{html_filename}")

    # 6. Move source to done/
    done_path = os.path.join(DONE_DIR, filename)
    shutil.move(file_path, done_path)
    print(f"  Source moved to: done/{filename}")

    print(f"  DONE!")
    return True


def main():
    """Main entry point — process all files in generate/."""
    print("=" * 60)
    print("  LECTURE SUMMARIZER")
    print("=" * 60)

    # Check for API key
    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("\nERROR: ANTHROPIC_API_KEY environment variable not set.")
        print("Set it with: export ANTHROPIC_API_KEY='your-key-here'")
        sys.exit(1)

    # Get course name from args or use default
    course_name = sys.argv[1] if len(sys.argv) > 1 else "מבני נתונים"
    university = sys.argv[2] if len(sys.argv) > 2 else "אוניברסיטת רייכמן"

    print(f"\nCourse: {course_name}")
    print(f"University: {university}")
    print(f"Looking in: {GENERATE_DIR}/")

    # Find all PDF/PPTX files in generate/
    files = []
    for ext in ["*.pdf", "*.PDF", "*.pptx", "*.PPTX"]:
        files.extend(glob.glob(os.path.join(GENERATE_DIR, ext)))

    if not files:
        print(f"\nNo PDF/PPTX files found in generate/")
        print(f"Drop your lecture files into: {GENERATE_DIR}/")
        sys.exit(0)

    print(f"\nFound {len(files)} file(s) to process:")
    for f in files:
        print(f"  - {os.path.basename(f)}")

    # Process each file
    success = 0
    failed = 0
    for file_path in sorted(files):
        try:
            if process_file(file_path, course_name, university):
                success += 1
            else:
                failed += 1
        except Exception as e:
            print(f"  ERROR processing {os.path.basename(file_path)}: {e}")
            failed += 1

    # Summary
    print(f"\n{'='*60}")
    print(f"  COMPLETE: {success} succeeded, {failed} failed")
    print(f"  Results in: {RESULTS_DIR}/")
    print(f"  Originals in: {DONE_DIR}/")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
